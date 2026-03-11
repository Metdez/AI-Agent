"""
Execution: PPTX text extraction including slide notes.
"""

import os


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract all text from a PPTX file, including slide notes.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX not found: {file_path}")

    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    try:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Notes: {notes_text}]")
            except Exception:
                pass

            if slide_parts:
                slides_text.append("\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)

        if len(full_text) < 50:
            raise ValueError(
                f"Could not extract text from PPTX. Only got {len(full_text)} characters. "
                "File may be image-only or corrupted."
            )

        return full_text

    except (ValueError, ImportError):
        raise
    except Exception as e:
        raise Exception(f"Failed to read PPTX '{file_path}': {e}")


if __name__ == "__main__":
    import sys
    from pathlib import Path
    sys.path.insert(0, str(Path(__file__).parent.parent))

    input_dir = Path(__file__).parent.parent / "input"
    pptx_files = list(input_dir.glob("*.pptx"))

    if not pptx_files:
        print("No PPTX files found in input/. Drop a PPTX there to test.")
    else:
        pptx = pptx_files[0]
        print(f"Reading: {pptx.name}")
        text = extract_text_from_pptx(str(pptx))
        print(f"Length: {len(text)} characters")
        print(f"First 500 chars:\n{text[:500]}")
